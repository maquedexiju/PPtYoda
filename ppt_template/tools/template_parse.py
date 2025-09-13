from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import re
import uuid
from lxml import etree

class Template_Parser():
    '''
    解析 ppt 模板文件
    1. 提取所有页面的模板
    2. 提取所有组件的模板
    '''

    def __init__(self, tempalte_file_path, logger, required_components = ['title', 'text']) -> None:

        # 设置 namespace
        self.ns = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
        }

        self.template_file_path = tempalte_file_path
        self.logger = logger
        self.prs = Presentation(self.template_file_path)

        self.sections = self.extract_sections()
        self.components, self.lost_components = self.extract_components(required_components)
        self.slide_templates = self.extract_slide_templates()


    def _get_slide(self, slide_id, ppt=None):

        if ppt is None: ppt = self.prs

        # 获取指定 slide
        if isinstance(slide_id, int):
            try:
                return ppt.slides[slide_id]
            except IndexError:
                pass

            try:
                return ppt.slides.get(slide_id)
            except IndexError:
                pass
        
        if isinstance(slide_id, str):

            try:
                slide_id = int(slide_id)
                try:
                    return ppt.slides[slide_id]
                except IndexError:
                    pass


                try:
                    return ppt.slides.get(slide_id)
                except IndexError:
                    pass
            except ValueError:
                pass

            for slide in ppt.slides:
                if slide.has_notes_slide:
                    slide_name = slide.notes_slide.notes_text_frame.text.split('\n')[0]
                    if slide_name == slide_id:
                        return slide
                if slide.name == slide_id:
                    return slide
            
            return None

        if isinstance(slide_id, Slide):
            return slide_id

    def _get_slide_name(self, slide):

        if slide.has_notes_slide:
            slide_name = slide.notes_slide.notes_text_frame.text.split('\n')[0]
            return slide_name

        # 返回 4 位 uuid
        return str(uuid.uuid4())[:4]

    def extract_sections(self):
        '''
        提取所有 section，返回格式：
        {
            'section_name': [
                {'id': slide_id, 'name': slide_name},
                ...
            ],
            ...
        }
        '''

        # 获取所有 section
        xml_str = self.prs.element.xml
        tree = etree.fromstring(xml_str)
        # xpath 找到 <p14:section name="section_name"> 中的所有 <p14:sldId>
        sections = tree.xpath(f"//p14:section", namespaces=self.ns)

        section_dict = {}
        for section in sections:
            section_name = section.get('name')
            if section_name == 'templates': continue
            section_dict[section_name] = [
                {
                    'id': slide_id,
                    'name': self._get_slide_name(self._get_slide(slide_id))
                } for slide_id in section.xpath("./p14:sldIdLst/p14:sldId/@id", namespaces=self.ns)
            ]

        return section_dict


    def extract_slide_templates(self):
        '''
        打开 template_file_path 的 ppt，遍历每一页
        1. 获取演讲者备注的信息，第一行为 名称，其他行为备注
        2. 获取页面中所有 @ 开头的 text，这些都是 placeholder，并获取 @ 后面的内容，冒号前为名称，冒号后为备注说明
        3. container 开头的元素是“容器”，容器中可以添加组件（#），括号里跟着样式描述，样式描述用空格隔开。示例：@container-toc(l-r l): 备注说明
            1. 展开方向 direction：lr 或 tb
            2. 水平对齐方式 align：l 或 c 或 r 或 j 
            3. 垂直对齐方式 valign：t 或 c 或 b
            4. 延伸方向的间隔 gap：数字
        返回一个 list，每个元素是一个 dict，包括页面序号、名称、备注、placeholder
        '''
        data_list = []
        for i, slide in enumerate(self.prs.slides):
            # 如果 slide_id 不在 templates 中，那么跳过
            if self.sections != {} and 'templates' in self.sections.keys() and slide.slide_id not in[int(slide['id']) for slide in  self.sections['templates']]:
                continue
            # 如果有章节，但没有 templates，那么跳过章节中的内容
            skip = False
            if self.sections != {} and 'templates' not in self.sections.keys():
                for section in self.sections:
                    if slide.slide_id in [int(slide['id']) for slide in self.sections[section]]:
                        skip = True
                        break
                if skip: continue
            
            slide_data = {
                'id': i,
                'name': '',
                'description': '',
                'placeholders': []
            }
            # 获取演讲者备注的信息
            if slide.notes_slide:
                note = slide.notes_slide.notes_text_frame.text
                note_list = note.split('\n')
                if len(note_list) > 0:
                    slide_data['name'] = note_list[0]
                    slide_data['description'] = '\n'.join(note_list[1:])
            
            # 如果没有名称，那么预警跳过
            if slide_data['name'] == '':
                # self.logger.warning(f"第 {i+1} 页没有名称，跳过")
                continue

            # # 获取页面中所有 @ 开头的 text
            # for shape in slide.shapes:
            #     if shape.has_text_frame:
            #         text_frame = shape.text_frame
            #         for paragraph in text_frame.paragraphs:
            #             if paragraph.text.startswith('@'):
            #                 print(paragraph.text)
            #                 placeholder = paragraph.text[1:]
            #                 # 描述可能不存在
            #                 placeholder_description = ''
            #                 if ':' in placeholder:
            #                     placeholder_name, placeholder_description = placeholder.split('：')
            #                 else:
            #                     placeholder_name = placeholder

            #                 placeholder_data = {
            #                     'name': placeholder_name,
            #                     'description': placeholder_description
            #                 }
            #                 slide_data['placeholder'].append(placeholder_data)
            placeholder_list = []
            for shape in slide.shapes:
                spl = self._get_placeholder(shape)
                if spl != []: placeholder_list.extend(spl)
            slide_data['placeholders'] = self._sort_placeholders(placeholder_list)

            data_list.append(slide_data)
        
        # self.template_data = data_list
        return data_list

    def _sort_placeholders(self, placeholder_list):
        '''
        对 placeholder 进行排序，按照 name 排序
        '''
        return sorted(placeholder_list, key=lambda x: x['name'])


    
    def _get_placeholder(self, shape):
        '''
        从 shape 中提取 placeholder，如果是 @-c 开头，那么是组件，这部分要通知上层的 group 进行跳过
        '''
        placeholder_list = []
        if shape.has_text_frame:
            text = shape.text_frame.text

            if text.startswith('#'):
                return 'component'

            if text.startswith('@'):

                text = text[1:]
                if text.startswith('img-'):
                    type = 'img'
                    text = text[4:]
                elif text.startswith('svg-'):
                    type = 'svg'
                    text = text[4:]
                elif text.startswith('container-'):
                    type = 'container'
                    text = text[10:]
                elif text.startswith('icon-'):
                    type = 'icon'
                    text = text[5:]
                else:
                    # diy 前缀需要保留在 name 中，同时 type 为空
                    type = 'text'

                # 寻找 name(style)：notes，(style) 和 ：notes 都可能不存在
                # pattern = r'^(?P<name>[^(:\n]+)(?:[\(（](?P<style>[^()]*)[\)）])?(?:[\:：](?P<notes>.*))?$'
                pattern = r'^(?P<name>[^(:：\(\)（）\n]+)(?:[(\uff08](?P<style>[^()（）]*?)[)\uff09])?(?:[:\uff1a](?P<notes>.*))?$'
                match = re.match(pattern, text)

                if match:
                    data = dict(name=match.group('name'))
                    data['type'] = type
                    if match.group('style') != None:
                        data['style'] = match.group('style')
                    if match.group('notes') != None:
                        data['description'] = match.group('notes')

                # 对 container 获取 component_placeholders min_component_number max_component_number
                if type == 'container':
                    component = [ c for c in self.components if c['name'] == data['name']]
                    if component != []:
                        data['component_placeholders'] = component[0]['placeholders']

                    style_dict = self._parse_style(data['style'])
                    data['min_component_number'] = int(style_dict['min_n']) if 'min_n' in style_dict else 1
                    data['max_component_number'] = int(style_dict['max_n']) if 'max_n' in style_dict else 5

                placeholder_list.append(data)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                sub_pl = self._get_placeholder(sub_shape)
                if sub_pl != 'component' and sub_pl != []: placeholder_list.extend(sub_pl)

        return placeholder_list

    
    def _parse_style(self, style_str):
        '''
        解析 style_str，返回一个 dict
        '''
        style_dict = {}
        style_list = style_str.split(' ')
        for style in style_list:
            if '-' in style:
                key, value = style.split('-', 1)
                style_dict[key] = value
            else:
                style_dict[style] = ''

        return style_dict

    
    def extract_components(self, required_components = ['title', 'text']):
        '''
        从页面中提取组件，组件是一个 group 元素，包含一个 # 开头的 text
        需要去除 # 开头的 text，把剩下元素的 xml 保存下来，作为元素，并根据 # 后跟随的文字命名
        返回 json
        '''
        component_list = []
        # 遍历所有编组的形状
        for slide in self.prs.slides:
            # 如果 slide_id 不在 templates 中，那么跳过
            if self.sections != {} and 'templates' in self.sections.keys() and slide.slide_id not in[int(slide['id']) for slide in  self.sections['templates']]:
                continue
            # 如果有章节，但没有 templates，那么跳过章节中的内容
            skip = False
            if self.sections != {} and 'templates' not in self.sections.keys():
                for section in self.sections:
                    if slide.slide_id in [int(slide['id']) for slide in self.sections[section]]:
                        skip = True
                        break
                if skip: continue

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP and len(shape.shapes) == 2:
                # 理论上只有两个元素，一个 text，另一个是 group、text 或其他
                    is_component = False
                    component_name = ''
                    component_xml = ''
                    sub_shapes = shape.shapes
                    sub_shape_1 = sub_shapes[0]
                    sub_shape_2 = sub_shapes[1]
                    if sub_shape_1.has_text_frame and sub_shape_1.text.startswith('#'):
                        is_component = True
                        component_name = sub_shape_1.text[1:]
                        component_xml = sub_shape_2.element.xml
                        component = sub_shape_2
                    elif sub_shape_2.has_text_frame and sub_shape_2.text.startswith('#'):
                        is_component = True
                        component_name = sub_shape_2.text[1:]
                        component_xml = sub_shape_1.element.xml
                        component = sub_shape_1

                    if is_component:
                        if component_name != '' and component_xml != '':
                            component_list.append({
                                'name': component_name,
                                # 'xml': component_xml,
                                'placeholders': self._sort_placeholders(self._get_placeholder(component))
                            })

        components_names = [item['name'] for item in component_list]
        lost_components = []
        for required_component in required_components:
            if required_component not in components_names:
                lost_components.append(required_component)

        return component_list, lost_components


    def export_templates(self, file_name):
        '''
        导出 slide_templates 到 file_name
        '''
        with open(file_name, 'w', encoding='utf-8') as f:
            json.dump(self.slide_templates, f, ensure_ascii=False, indent=4)


    def export_components(self, file_name):
        '''
        导出 components 到 file_name
        '''
        with open(file_name, 'w', encoding='utf-8') as f:
            json.dump(self.components, f, ensure_ascii=False, indent=4)



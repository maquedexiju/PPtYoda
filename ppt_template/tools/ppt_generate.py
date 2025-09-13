from pydoc import text
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.slide import Slide
from pptx.opc.package import Part
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
import os
from lxml import etree
import shutil
import copy
import tempfile
import zipfile
import uuid
import re
from openai import AsyncOpenAI, BadRequestError
import base64
from PIL import Image
import io
import asyncio

from pptx.oxml.ns import _nsmap

_nsmap['p14'] = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
_nsmap['asvg'] = 'http://schemas.microsoft.com/office/drawing/2016/SVG/main'

from ..models import PPt_Template
from .ppt_color_parse import Color_Parser
from .icons_handle import Icons_Handler

from django.conf import settings

# 获取 __file__ 的绝对路径
file_path = os.path.abspath(__file__)
# 获取 __file__ 父目录的父目录
file_dir = os.path.dirname(os.path.dirname(file_path))
# 获取父目录的 templates 目录
templates_dir = os.path.join(file_dir, 'templates')

default_img_prompt = '简约插画风' # '请不要在图片中增加文字'
default_svg_prompt = ''

# 从路径模板新建 ppt
class PPt_Generator:
    '''
    根据 template 和生成的 slide_date 生成 ppt
    目前还添加了很多插入现有页面的功能，后续可能有用，或者分拆出来
    '''

    def __init__(self, file_name: str, template: PPt_Template):

        # 设置 namespace
        self.ns = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        }

        # 处理 template
        self.template = template
        self.template_path = template.file.path
        self.temp_ppt = Presentation(self.template_path)
        self.color_parser = Color_Parser(self.temp_ppt)

        self.temp_dir = tempfile.TemporaryDirectory()
        self.file_name = file_name


        if not os.path.exists(file_name):

            # 复制模板到临时文件
            tmp_file_name = os.path.join(self.temp_dir.name, 'temp.pptx')
            shutil.copy(self.template_path, tmp_file_name)
            # self.tmp_file_path = tmp_file_name
            self.ppt = Presentation(tmp_file_name)
            # 删除 ppt 中的所有页面
            # for slide in self.ppt.slides:
            #     self.ppt.slides.remove(slide)
            # xml_slides = self.ppt.slides._sldIdLst
            # for slide in xml_slides:
            #     rId = slide.rId
            #     self.ppt.part.drop_rel(rId)
            #     del slide
                # xml_slides.remove(slide)
            for slide in self.ppt.slides:
                self._rm_slide(slide)
            self._rm_sections()

        else:
            self.ppt = Presentation(file_name)

        # self.component_list = template.components
        self.component_list = self._extract_components()
        self.section_dict = self._get_sections()

        # llm 设置
        self.client = AsyncOpenAI(base_url=settings.LLM_IMG_GEN_BASE_URL, api_key=settings.LLM_IMG_GEN_API_KEY)
        self.img_gen_model = settings.LLM_IMG_GEN_MODEL

        self.icons_handler = Icons_Handler()

    def _get_sections(self):

        # 获取所有 section
        xml_str = self.ppt.element.xml
        tree = etree.fromstring(xml_str)
        # xpath 找到 <p14:section name="section_name"> 中的所有 <p14:sldId>
        sections = tree.xpath(f"//p14:section", namespaces=self.ns)

        section_dict = {}
        for section in sections:
            section_name = section.get('name')
            section_dict[section_name] = section.xpath("./p14:sldIdLst/p14:sldId/@id", namespaces=self.ns)

        return section_dict

    def _rm_sections(self):

        # 获取所有 section
        # sections = self.ppt._element.xpath(f"//p14:section")
        # section_list = self.ppt._element.xpath(f"//p14:sectionLst")
        # if len(section_list) == 0:
        #     return

        # section_list = section_list[0]
        # for section in sections:
        #     section_list.remove(section)

        section_list = self.ppt._element.xpath(f"//p14:sectionLst")
        p_ext = self.ppt._element.xpath('//p:ext')
        if section_list:
            p_ext[0].remove(section_list[0])

    def _get_slide(self, slide_id, ppt=None):

        if ppt is None: ppt = self.ppt

        # 获取指定 slide
        if isinstance(slide_id, int):
            try:
                return ppt.slides[slide_id]
            except IndexError:
                pass

            try:
                return ppt.slides.get(slide_id)
            except IndexError:
                pass
        
        if isinstance(slide_id, str):

            try:
                slide_id = int(slide_id)
                try:
                    return ppt.slides[slide_id]
                except IndexError:
                    pass


                try:
                    return ppt.slides.get(slide_id)
                except IndexError:
                    pass
            except ValueError:
                pass

            for slide in ppt.slides:
                if slide.has_notes_slide:
                    slide_name = slide.notes_slide.notes_text_frame.text.split('\n')[0]
                    if slide_name == slide_id:
                        return slide
                if slide.name == slide_id:
                    return slide
            
            return None

        if isinstance(slide_id, Slide):
            return slide_id


    def _rm_slide(self, slide_id):

        # 删除指定 slide

        slide = self._get_slide(slide_id)
        slide_id = slide.slide_id

        xml_slides = self.ppt.slides._sldIdLst
        rel_list = [x for x in xml_slides if x.id==slide_id]
        rel = rel_list[0]
        rel_id = rel.rId

        xml_slides.remove(rel) # 删除页面列表中的索引
        self.ppt.part.drop_rel(rel_id) # 是啥对应的 xml
        # 保存时会自动删除没有索引的素材


    def save(self, file_name_with_path=None):

        # self.ppt.save(self.tmp_file_path)
        # # 解压文件
        # with zipfile.ZipFile(self.tmp_file_path, 'r') as zip_ref:
        #     unzip_dir = os.path.join(self.temp_dir.name, 'unzip')
        #     zip_ref.extractall(unzip_dir)
        

        # 复制模板到指定路径
        if file_name_with_path is None:
            file_path = os.path.dirname(self.file_name)
            file_name = os.path.basename(self.file_name)
        else:
            file_path = os.path.dirname(file_name_with_path)
            file_name = os.path.basename(file_name_with_path)

        if file_path != '' and not os.path.exists(file_path):
            os.makedirs(file_path)

        self.ppt.save(os.path.join(file_path, file_name))

        # with zipfile.ZipFile(self.file_name, 'r') as zip_ref:
        #     zip_ref.extractall(self.file_name[:-5])

    def save_and_export_blob(self, file_name_with_path=None):
        '''
        导出 ppt 文件为 blob
        '''
        self.save(file_name_with_path)
        if file_name_with_path is None:
            file_name_with_path = self.file_name
        with open(file_name_with_path, 'rb') as f:
            blob = f.read()
        return blob

    def __del__(self):
        self.temp_dir.cleanup()

# ============
# 模板处理相关
# ============

    def _calc_group_shape_size(self, group_shape):
        '''
        计算 group 形状的大小
        返回：
            width, height
        '''

        left, top, right, bottom = Emu(0), Emu(0), Emu(0), Emu(0)
        width, height = [], []

        for shape in group_shape.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                g_width, g_height = self._calc_group_shape_size(shape)
                width.append(g_width)
                height.append(g_height)
            else:
                if shape.left < left:
                    left = shape.left
                if shape.top < top:
                    top = shape.top
                if shape.left + shape.width > right:
                    right = shape.left + shape.width
                if shape.top + shape.height > bottom:
                    bottom = shape.top + shape.height
        
        width.append(right - left)
        height.append(bottom - top)
        # 取最大的宽度和高度
        width = max(width)
        height = max(height)

        return width, height

    def _extract_components(self):
        '''
        从页面中提取组件，组件是一个 group 元素，包含一个 # 开头的 text
        需要去除 # 开头的 text，把剩下元素的 xml 保存下来，作为元素，并根据 # 后跟随的文字命名
        返回 json
        '''
        component_list = []
        # 遍历所有编组的形状
        for slide in self.temp_ppt.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP and len(shape.shapes) == 2:
                # 理论上只有两个元素，一个 text，另一个是 group、text 或其他
                    is_component = False
                    component_name = ''
                    component_element = None
                    sub_shapes = shape.shapes
                    sub_shape_1 = sub_shapes[0]
                    sub_shape_2 = sub_shapes[1]
                    if sub_shape_1.has_text_frame and sub_shape_1.text.startswith('#'):
                        is_component = True
                        component_name = sub_shape_1.text[1:]
                        # component_element = sub_shape_2._element
                        component = sub_shape_2
                    elif sub_shape_2.has_text_frame and sub_shape_2.text.startswith('#'):
                        is_component = True
                        component_name = sub_shape_2.text[1:]
                        # component_element = sub_shape_1._element
                        component = sub_shape_1

                    if is_component:

                        # 计算比例尺
                        origin_width = int(shape.element.xpath('./p:grpSpPr/a:xfrm/a:ext/@cx')[0])
                        scaled_width = int(shape.element.xpath('./p:grpSpPr/a:xfrm/a:chExt/@cx')[0])
                        scale = origin_width / scaled_width


                        # 获取 tag 的 rid
                        # tags = {}
                        # rids = component._element.xpath('.//p:tags/@r:id')
                        # for rid in rids:
                        #     r_part = component.part.related_part(rid)
                        #     tags[rid] = r_part

                        # if component.shape_type == MSO_SHAPE_TYPE.GROUP:
                        #     width, height = self._calc_group_shape_size(component)
                        # else:
                        #     width = component.width
                        #     height = component.height
                        component_list.append({
                            'name': component_name,
                            # 'xml': component_xml,
                            'element': component,
                            'width': component.width * scale,
                            'height': component.height * scale,
                            # 'tags': tags,
                            'placeholders': self._get_placeholder(component)
                        })

        return component_list

    
    def _get_placeholder(self, shape):
        '''
        从 shape 中提取 placeholder，如果是 @-c 开头，那么是组件，这部分要通知上层的 group 进行跳过
        '''
        placeholder_list = []
        if shape.has_text_frame:
            text = shape.text_frame.text

            if text.startswith('#'):
                return 'component'

            if text.startswith('@'):

                text = text[1:]
                if text.startswith('img-'):
                    type = 'img'
                    text = text[4:]
                elif text.startswith('svg-'):
                    type = 'svg'
                    text = text[4:]
                elif text.startswith('container-'):
                    type = 'container'
                    text = text[10:]
                else:
                    type = 'text'

                # 寻找 name(style)：notes，(style) 和 ：notes 都可能不存在
                pattern = r'^(?P<name>[^(:\n]+)(?:[\(（](?P<style>[^()]*)[\)）])?(?:[\:：](?P<notes>.*))?$'
                match = re.match(pattern, text)

                if match:
                    data = dict(name=match.group('name'))
                    data['type'] = type
                    if match.group('style') != None:
                        data['style'] = match.group('style')
                    if match.group('notes') != None:
                        data['description'] = match.group('notes')

                # 对 container 获取 component_placeholders min_component_number max_component_number
                if type == 'container':
                    component = [ c for c in self.components if c['name'] == data['name']]
                    if component != []:
                        data['component_placeholders'] = component[0]['placeholders']

                    style_dict = self._parse_style(data['style'])
                    data['min_component_number'] = int(style_dict['min_n']) if 'min_n' in style_dict else 1
                    data['max_component_number'] = int(style_dict['max_n']) if 'max_n' in style_dict else 5

                placeholder_list.append(data)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                sub_pl = self._get_placeholder(sub_shape)
                if sub_pl != 'component' and sub_pl != []: placeholder_list.extend(sub_pl)

        return placeholder_list

    
    def _get_slides_id_from_temp_section(self, section_name):
        # 获取指定 section 中的所有 slide
        xml_str = self.ppt.element.xml
        tree = etree.fromstring(xml_str)
        # xpath 找到 <p14:section name="section_name"> 中的所有 <p14:sldId>
        slide_ids = tree.xpath(f"//p14:section[@name='{section_name}']//p14:sldId/@id", namespaces=self.ns)

        return slide_ids

    def _save_temp_img(self, img):

        name = img.filename
        dpi = img.dpi
        size = img.size
        blob = img.blob
        content_type = img.content_type
        # 保存图片到临时文件
        temp_file_path = os.path.join(self.temp_dir.name, name)
        with open(temp_file_path, 'wb') as f:
            f.write(blob)

        return temp_file_path

    def _create_image_part(self, n_shps, n_shp, o_shp):
        # 在新的 shapes 下，把 o_shp 中的图片复制过来

        temp_img_file = self._save_temp_img(o_shp.image)
        img_part, rId = n_shps.part.get_or_add_image_part(temp_img_file)
        bilp_info = n_shp._element.xpath('./p:blipFill/a:blip')[0]
        bilp_info.set(f'{{{self.ns['r']}}}embed', rId)

        # 判断是否是 svg
        svg_embed = o_shp._element.xpath('p:blipFill/a:blip//asvg:svgBlip/@r:embed')
        if svg_embed != []:
            rId = svg_embed[0]
            svg_part =o_shp.part.related_part(rId)
            n_package = n_shps.part._package
            n_svg_part = Part(n_package.next_image_partname('svg'), 'image/svg+xml', n_package, svg_part.blob)
            # n_svg_part = Part(n_package.next_media_partname('svg'), 'image/svg+xml', n_package, svg_part.blob)
            n_rId = n_shps.part.relate_to(n_svg_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')

            n_svg_blip = n_shp._element.xpath('p:blipFill/a:blip//asvg:svgBlip')[0]
            n_svg_blip.set(f'{{{self.ns['r']}}}embed', n_rId)

    
    def _create_image_parts(self, n_shps, o_shps):
        # 创建图片文件
        
        for shp in n_shps:

            shp_id = shp.shape_id
            o_shp = [ s for s in o_shps if s.shape_id == shp_id][0]

            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                self._create_image_part(n_shps, shp, o_shp)

            elif shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._create_image_parts(shp.shapes, o_shp.shapes)


    def _deepcopy_shapes(self, n_shps, o_shps):

        for shp in o_shps:

            if shp.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                continue

            # 复制 shape
            el = shp.element
            el_n = copy.deepcopy(el)
            n_shps._spTree.insert_element_before(el_n)
            n_shp = n_shps[-1]
            
            # 处理图片
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                self._create_image_part(n_shps, n_shp, shp)
            # 处理分组下的图片
            elif shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._create_image_parts(n_shp.shapes, shp.shapes)

            # 处理 tags
            # rids = shp._element.xpath('.//p:tags/@r:id')
            # for rid in rids:
            #     r_part = shp.part.related_part(rid)
            #     n_rid = n_shp.part.relate_to(r_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags")
            #     p_tags = n_shp._element.xpath(f'.//p:tags[@r:id="{rid}"]')
            #     p_tags[0].set(f'{{{self.ns['r']}}}id', n_rid)
                    

    def _get_slide_layout_index(self, prs, slide):

        i = 0
        layout = slide.slide_layout
        for slide_master in prs.slide_masters:
            j = 0
            for slide_layout in slide_master.slide_layouts:
                if slide_layout == layout:
                    return i, j
                j += 1
            i += 1


    def _insert_slide_from_template(self, slide_id, copy_notes=False):
        # 复制模板中的指定 slide 到 ppt
        slide = self._get_slide(slide_id, self.temp_ppt)
        mi, si = self._get_slide_layout_index(self.temp_ppt, slide)
        slide_n = self.ppt.slides.add_slide(self.ppt.slide_masters[mi].slide_layouts[si])

        # slide_n._element = copy.deepcopy(slide._element)

        # return slide_n

        # 如果背景是单独的
        if not slide.follow_master_background:
            slide_n.background._element = copy.deepcopy(slide.background._element)

        # 再处理 placeholders
        for placeholder in slide.placeholders:

            o_text_frame = placeholder.text_frame._element
            idx = placeholder._element.ph_idx
            try:
                n_text_frame = slide_n.placeholders[idx].text_frame._element
                slide_n.placeholders[idx]._element.replace(n_text_frame, copy.deepcopy(o_text_frame))
            except: # KeyError("no placeholder on this slide with idx == %d" % idx)
                pass

            # n_text_frame = slide_n.placeholders[i].text_frame._element
            # o_text_frame = slide.placeholders[i].text_frame._element

        self._deepcopy_shapes(slide_n.shapes, slide.shapes)

        if copy_notes:
            # 复制 notes
            notes = slide.notes_slide
            notes_n = slide_n.notes_slide
            notes_n.notes_text_frame.text = notes.notes_text_frame.text
        # for shp in slide.shapes:

        #     if shp.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        #         continue
            
        #     # 处理图片
        #     if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
        #         # 复制图片
        #         image = shp.image
        #         temp_img_file = self._save_temp_img(image)
        #         image_n = slide_n.shapes.add_picture(temp_img_file, shp.left, shp.top, shp.width, shp.height)
        #         # image_n = slide_n.shapes.add_picture(image.image, image.left, image.top, image.width, image.height)
        #         # 复制图片的属性
        #         image_n._element = copy.deepcopy(shp._element)
        #     else:
        #         # 复制 shape
        #         el = shp.element
        #         newel = copy.deepcopy(el)
        #         slide_n.shapes._spTree.insert_element_before(newel, 'p:extLst')

        return slide_n


    def add_slides(self, slides):
        # 添加 slides 到 ppt
        for section_name in slides.keys():

            if not section_name.startswith('unnamed') and section_name.startswith('@'):
                # 复制模板中的指定 section 到 ppt
                section_name = section_name[1:]
                slides_name = self._get_slides_id_from_temp_section(section_name)
                for slide_name in slides_name:
                    self._insert_slide_from_template(slide_name, copy_notes=True)


# ============
# llm
# ============

    def get_img_size(self, width, height):
        # 如果 width 不是 int 嚯 float，获取 pt 属性
        if isinstance(width, Emu):
            width = width.pt
        if isinstance(height, Emu):
            height = height.pt

        if height*width < 262144:
            ratio = height*width/262144
            height = int(height/ratio)
            width = int(width/ratio)
        elif height*width > 4194304:
            ratio = height*width/4194304
            height = int(height/ratio)
            width = int(width/ratio)
        else:
            ratio = 1

        if width > 2048 and width > height:
            width = 2048
            height = int(height*2048/width)
        if height > 2048 and height > width:
            height = 2048
            width = int(width*2048/height)

        return width, height


    async def _generate_img(self, uid, desc, width, height):
        '''
        通过 AI 生成图片，返回值：是否成功、成功时返回图片路径，失败时返回错误信息
        '''

        width, height = self.get_img_size(width, height)

        # 计算长宽比
        aspect_ratio = width/height
        if aspect_ratio > 16 or aspect_ratio < 1/16:
            return False

        size = f'{int(width)}x{int(height)}'
        try:
            response = await self.client.images.generate(
                prompt=desc,
                response_format="b64_json",  # 关键：返回 Base64 而非 URL
                size=size,
                model=self.img_gen_model,
                extra_body={
                    'watermark': False
                }
            )
        except BadRequestError as e:
            return False, e.response.json()['error']['message']
        
        image_data = base64.b64decode(response.data[0].b64_json)
        img_path = os.path.join(self.temp_dir.name, f'{uid}.jpg')
        with open(img_path, "wb") as f:
            f.write(image_data)
        return True, img_path


# ============
# 通过 slide data 来生成 PPt
# ============

    def add_slide(self, slide_data, slide_note=None):

        if type(slide_data) == dict and 'template_id' in slide_data.keys():
        # 从数据生成页面
            temp_no = int(slide_data['template_id'])
            new_slide = self._insert_slide_from_template(temp_no)
            slide_notes = self.replace_placeholders(new_slide.shapes, slide_data['placeholders'])
            if slide_note:
                slide_notes.extend([slide_note])

            new_slide.notes_slide.notes_text_frame.text = '\n'.join(slide_notes)
            self.replace_icons(new_slide.shapes)

            # 去掉 p:nvPr 下的 p:custDataLst，还是很容易出错
            # 尤其是在 container 中时，可能 rid 不能出现多次吧
            for nvPr in new_slide._element.xpath('//p:nvPr'):
                if cust := nvPr.xpath('./p:custDataLst'):
                    for c in cust:
                        nvPr.remove(c)
        
        else:
        # 直接添加一系列的 slide，添加模板页
            # for section_name in slide_data.keys():
            #     for slide_id in slide_data[section_name]:
            #         slide = self._get_slide(slide_id, self.temp_ppt)
            #         new_slide = self._insert_slide_from_template(slide)
            for slide_id in slide_data:
                slide = self._get_slide(slide_id, self.temp_ppt)
                new_slide = self._insert_slide_from_template(slide, copy_notes=True)
            

        return new_slide


    def _replace_text(self, shape, text):

        shape_text_frame = shape.text_frame
        first_para = shape_text_frame.paragraphs[0]

        # 删除除第一个 run
        for r in first_para.runs[1:]:
            first_para._element.remove(r._r)

        template_p = first_para
        shape_text_frame._element.remove_all('a:p')

        # 经过上述步骤，获取了模板样式，删除了所有段落
        text_list = text.split('\n')
        if len(text_list) == 1 and '；' in text_list[0]:
            text_list = text_list[0].split('；')
        for line in text_list:
            tmp_p = copy.deepcopy(template_p)
            tmp_p.runs[0].text = line
            shape_text_frame._element.insert(-1, tmp_p._element)


    def _replace_icon_shape(self, shape, shapes, icon_name, ratio=1):

        icon_path = self.icons_handler.find_icon(icon_name)
        if icon_path:
            icon_color = self.color_parser.get_fore_color_rgb_hex(shape)
            width = int(shape.width.pt * ratio)
            height = int(shape.height.pt * ratio)
            icon_svg_content = self.icons_handler.get_icon_content(icon_path)
            icon_png_path = os.path.join(self.temp_dir.name, f'{icon_name}.png')
            icon_svg_path = os.path.join(self.temp_dir.name, f'{icon_name}.svg')
            modified_icon_svg_content = self.icons_handler.modify_svg(icon_svg_content, icon_svg_path, icon_color)
            self.icons_handler.svg_to_png(modified_icon_svg_content, icon_png_path, width, height)
            pic = shapes.add_picture(icon_png_path, shape.left, shape.top, shape.width, shape.height)
            shapes._element.remove(shape._element)

            package = pic.part._package
            svg_part = Part(package.next_image_partname('svg'), 'image/svg+xml', package, modified_icon_svg_content.encode('utf8'))
            rid = pic.part.relate_to(svg_part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/image')
            extLst_str = '''
            <a:extLst>
                <a:ext uri="{96DAC541-7B7A-43D3-8B79-37D633B846F1}">
                    <asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main"
                        r:embed="RID" />
                </a:ext>
            </a:extLst>
            '''.replace('RID', rid)
            parser = etree.XMLParser(recover=True)
            pic._element.xpath('p:blipFill/a:blip')[0].append(etree.fromstring(extLst_str, parser))
        else:
            self._replace_text(shape, f'@icon-{icon_name}')
            

    def _parse_style(self, style_str):
        '''
        解析 style_str，返回一个 dict
        '''
        style_dict = {}
        style_list = style_str.split(' ')
        for style in style_list:
            if '-' in style:
                key, value = style.split('-', 1)
                style_dict[key] = value
            else:
                style_dict[style] = ''

        return style_dict

    def _calculate_element_positions(
        self,
        parent_width: Emu, parent_height: Emu, parent_x: Emu, parent_y: Emu,
        child_width: Emu, child_height: Emu,
        child_count: int,
        direction: str = "l-r",
        horizontal_alignment: str = "l",
        vertical_alignment: str = "t",
        spacing: Emu = Emu(0)
    ) -> list[tuple[float, float]]:
        """
        计算子元素在父元素中的排列位置
        
        参数:
            parent_width: 父元素宽度
            parent_height: 父元素高度
            parent_x: 父元素左上角x坐标
            parent_y: 父元素左上角y坐标
            child_width: 子元素宽度
            child_height: 子元素高度
            child_count: 子元素数量
            direction: 排列方式，"l-r" 表示从左到右，"t-b" 表示从上到下
            horizontal_alignment: 水平对齐方式，可选值："l", "r", "c", "j"
            vertical_alignment: 垂直对齐方式，可选值："t", "b", "c", "j"
            spacing: 子元素之间的间隔大小
        
        返回:
            每个子元素的左上角坐标列表，格式为 [(x1, y1), (x2, y2), ...]
        """
        positions = []
        
        # 计算子元素总体尺寸
        if direction == "l-r":
            total_children_width = child_width * child_count + spacing * (child_count - 1)
            total_children_height = child_height
        else:  # t-b
            total_children_width = child_width
            total_children_height = child_height * child_count + spacing * (child_count - 1)
        
        # 计算水平方向偏移量
        if horizontal_alignment == "l":
            start_x = parent_x
        elif horizontal_alignment == "r":
            start_x = parent_x + parent_width - total_children_width
        elif horizontal_alignment == "c":
            start_x = parent_x + (parent_width - total_children_width) / 2
        else:  # j (只在l-r排列时有效)
            if child_count > 1:
                # 重新计算 j 模式下的间距
                spacing = (parent_width - child_width * child_count) / (child_count - 1)
            start_x = parent_x
        
        # 计算垂直方向偏移量
        if vertical_alignment == "t":
            start_y = parent_y
        elif vertical_alignment == "b":
            start_y = parent_y + parent_height - total_children_height
        elif vertical_alignment == "c":
            start_y = parent_y + (parent_height - total_children_height) / 2
        else:  # j (只在t-b排列时有效)
            if child_count > 1:
                # 重新计算 j 模式下的间距
                spacing = (parent_height - child_height * child_count) / (child_count - 1)
            start_y = parent_y
        
        # 计算每个子元素的位置
        for i in range(child_count):
            if direction == "l-r":
                # 从左到右排列
                x = start_x + i * (child_width + spacing)
                y = start_y
            else:  # t-b
                # 从上到下排列
                x = start_x
                y = start_y + i * (child_height + spacing)
            
            positions.append((Emu(x), Emu(y)))
        
        return positions

    def _replace_container(self, shape, placeholder):
        '''
        替换 container
        args:
            shape: python-pptx 中的 shape，shape 中的文字包含了样式信息
            placeholder: 一个 dict，key 包括 name、components_placeholders
        '''

        # 获取 component
        component_name = placeholder['name']
        component_dict = [c for c in self.component_list if c['name'] == component_name]
        component = component_dict[0]['element']

        # 解析 style
        ## 提供默认 style
        style = {
            'direction': 'l-r',
            'align': 'c',
            'valign': 'c',
            'gap': '160000'
        }
        shape_text = shape.text_frame.text
        # 寻找 name(style)：notes，(style) 和 ：notes 都可能不存在
        pattern = r'^(?P<name>[^(:\n]+)(?:[\(（](?P<style>[^()]*)[\)）])?(?:[\:：](?P<notes>.*))?$'
        match = re.match(pattern, shape_text)
        if match and match.group('style'):
            mod_style_str = match.group('style')
            for k, v in self._parse_style(mod_style_str).items():
                style[k] = v

        # 获取 component 数量
        components_len = len(placeholder['components_placeholders'])

        # 获取 shape 的宽、高、左、上
        width = shape.width
        height = shape.height
        left = shape.left
        top = shape.top

        # 获取 component 的宽、高
        component_width = component_dict[0]['width']
        component_height = component_dict[0]['height']

        # 计算子元素位置
        positions = self._calculate_element_positions(
            width, height, left, top,
            child_width=component_width, child_height=component_height,
            child_count=components_len,
            direction=style['direction'],
            horizontal_alignment=style['align'],
            vertical_alignment=style['valign'],
            spacing=Emu(int(style['gap']))
        )

        # 创建 component，然后添加到 ppt 上
        for i, pos in enumerate(positions):
            comp = copy.deepcopy(component)
            comp.left = pos[0]
            comp.top = pos[1]
            comp.width = int(component_width)
            comp.height = int(component_height)
            self.replace_placeholders([comp], placeholder['components_placeholders'][i])
            shape.element.addprevious(comp.element)

            # 添加 tags
            # tags = component_dict[0]['tags']

            # for rid in tags.keys():
            #     r_part = tags[rid]
            #     n_rid = comp.part.relate_to(r_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags")
            #     xpath_exp = f'.//p:tags[@r:id="{rid}"]'
            #     p_tags = comp._element.xpath(xpath_exp)
            #     p_tags[0].set(f'{{{self.ns['r']}}}id', n_rid)

        # 删除 shape
        shape.part.slide.shapes.element.remove(shape.element)


    def replace_placeholders(self, shapes, placeholders: list, ratio=1):
        '''
        遍历 shapes，如果 has_text_frame，读取 text_frame 信息，如果和 placeholder 匹配，替换其内容
        args:
            shapes: python-pptx 中的 shapes
            placeholders: 一个列表，每个元素是一个 dict，key 包括 name、content、type
        '''

        # 这里只是把文字描述添加到 ppt 中，不做 multimedia 的生成
        # multimedia 生成统一完成，不然针对 container 没有添加到 ppt 中，因为没有父元素，会报错：
        # 'NoneType' object has no attribute 'recalculate_extents'

        slide_notes = []

        for shape in shapes:
            # 判断 shape 是否是 group
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:

                origin_width = int(shape.element.xpath('./p:grpSpPr/a:xfrm/a:ext/@cx')[0])
                scaled_width = int(shape.element.xpath('./p:grpSpPr/a:xfrm/a:chExt/@cx')[0])
                scale = origin_width / scaled_width
                group_ratio = ratio * scale
                sub_slide_notes = self.replace_placeholders(shape.shapes, placeholders, group_ratio)
                slide_notes.extend(sub_slide_notes)
                continue

            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text_frame.text
                # 判断是否是 placeholder
                if not text.startswith('@'): continue

                # 寻找 type-name(style)：notes，type- (style) 和 ：notes 都可能不存在
                pattern = re.compile(
                    r'^(?:(?P<type>\w+)-)?'  # type部分变为可选（带连字符）
                    r'(?P<name>[^(:：\(\)（）\n]+)'
                    r'(?:[(\uff08](?P<style>[^()（）]*?)[)\uff09])?'
                    r'(?:[:\uff1a](?P<notes>.*))?$'
                )
                match = re.match(pattern, text[1:])
                if not match: continue
                name = match.group('name')

                placeholder = [p for p in placeholders if p['name'].replace(' ', '') == name]

                if len(placeholder) == 0: continue
                placeholder = placeholder[0]

                # 替换内容
                t = placeholder['type']

                if t == 'text':
                    # shape.text_frame.text = placeholder['content']
                    self._replace_text(shape, placeholder['content'])
                elif t == 'icon':
                    # self._replace_icon_shape(shape, shapes, placeholder['content'], ratio)
                    self._replace_text(shape, f'@icon-{placeholder["content"]}')

                elif t == 'img':
                    # 设置 uid 为 uuid 前四位
                    uid = str(uuid.uuid4())[:4]
                    # shape.text_frame.text = '@img-' + uid
                    t = '@img-' + uid
                    self._replace_text(shape, t)
                    # 在 slide_notes 中添加 uid 对应的图片
                    width, height = self.get_img_size(shape.width*ratio, shape.height*ratio)
                    slide_notes.append(
                        # f'@img-{uid}\nsize: {int(width)}*{int(height)}\n{default_img_prompt}\n{placeholder["content"]}\n@endimg',
                        f'@img-{uid}\n{default_img_prompt}\n{placeholder["content"]}\n@endimg',
                    )
                # elif t == 'svg':
                #     # 设置 uid 为 uuid 前四位
                #     uid = str(uuid.uuid4())[:4]
                #     t = '@svg-' + uid
                #     self._replace_text(shape, t)
                #     # 在 slide_notes 中添加 uid 对应的 svg
                #     slide_notes.append(
                #         f'@svg-{uid}\n{default_svg_prompt}\n{placeholder["content"]}\n@endsvg',
                #     )
                elif t == 'container':
                    self._replace_container(shape, placeholder)
        
        return slide_notes

    
    def replace_icons(self, shapes, ratio=1):
        '''
        遍历 shapes，如果 has_text_frame，读取 text_frame 信息，如果和 placeholder 匹配，替换其内容
        args:
            shapes: python-pptx 中的 shapes
        '''

        # 这里只是把文字描述添加到 ppt 中，不做 multimedia 的生成
        # multimedia 生成统一完成，不然针对 container 没有添加到 ppt 中，因为没有父元素，会报错：
        # 'NoneType' object has no attribute 'recalculate_extents'

        for shape in shapes:
            # 判断 shape 是否是 group
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:

                origin_width = int(shape.element.xpath('./p:grpSpPr/a:xfrm/a:ext/@cx')[0])
                scaled_width = int(shape.element.xpath('./p:grpSpPr/a:xfrm/a:chExt/@cx')[0])
                scale = origin_width / scaled_width
                group_ratio = ratio * scale
                group_top = shape.top
                group_left = shape.left
                group_width = shape.width
                group_height = shape.height
                self.replace_icons(shape.shapes, group_ratio)
                shape.top = group_top
                shape.left = group_left
                shape.width = group_width
                shape.height = group_height
                continue

            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text_frame.text
                # 判断是否是 placeholder
                if not text.startswith('@'): continue

                # 寻找 type-name(style)：notes，type- (style) 和 ：notes 都可能不存在
                pattern = re.compile(
                    r'^icon-(.*)'
                )
                match = re.match(pattern, text[1:])
                if not match: continue
                name = match.group(1)
                self._replace_icon_shape(shape, shapes, name, ratio)

# ============
# 图片替换与生成相关
# ============

    def add_colorful_text(self, shape, text):
        '''
        为 shape 中的文字添加颜色
        args:
            shape: python-pptx 中的 shape
            text: 要添加颜色的文字
        '''
        p = shape.text_frame.paragraphs[0]
        p.text = ''
        color_sheme = [
            # 好看的、高饱和度、底色度的彩虹色
            RGBColor(255, 0, 0),
            # RGBColor(255, 127, 0),
            # RGBColor(255, 255, 0),
            # RGBColor(127, 255, 0),
            # RGBColor(0, 255, 0),
            # RGBColor(0, 255, 127),
            # RGBColor(0, 255, 255),
            # RGBColor(0, 127, 255),
            # RGBColor(0, 0, 255),
            # RGBColor(127, 0, 255),
            # RGBColor(255, 0, 255),
            # RGBColor(255, 0, 127),
        ]
        c_i = 0
        for t in text:
            r = p.add_run()
            r.text = t
            r.font.color.rgb = color_sheme[c_i]
            c_i += 1
            if c_i >= len(color_sheme):
                c_i = 0

    async def replace_shapes_img(self, shapes, slide_notes_str):
        '''
        替换 shapes 中的图片
        args:
            shapes: python-pptx 中的 shapes
        '''

        for shp in shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                await self.replace_shapes_img(shp.shapes, slide_notes_str)
            elif hasattr(shp, 'has_text_frame') and shp.has_text_frame and shp.text_frame.text.startswith('@img'):
                text = shp.text_frame.text
                match = re.search(r'@img-(.+?)(-g)?$', text)
                if match and match.group(2) is None:
                    # 获取信息
                    uid = match.group(1)
                    reg = re.compile(f'@img-{uid}\n(.*?)\n@endimg', re.DOTALL)
                    try:
                        img_desc = re.search(reg, slide_notes_str).group(1)
                    except:
                        continue
                    # 移除之前生成的图片
                    for p in shapes:
                        if p.name == f'img-{uid}':
                            shapes._spTree.remove(p.element)
                            break
                    
                    # 生成图片
                    suc, file_path = await self._generate_img(uid, img_desc, shp.width, shp.height)
                    if not suc:
                        # 添加错误信息到材料中
                        new_text = text + ' ' + file_path
                        self.add_colorful_text(shp, new_text)
                    else:
                        # 替换图片
                        p = shapes.add_picture(file_path, width = shp.width, height = shp.height, left = shp.left, top = shp.top)
                        shapes._spTree.remove(p.element)
                        p.name = f'img-{uid}'
                        shp._element.addprevious(p.element)
                        # 文字替换为 -g
                        new_text = text.replace(uid, uid + '-g')
                        self.add_colorful_text(shp, new_text)


    async def replace_slide_img(self, slide):
        '''
        获取 slide 中 @img 开头的 shape，并获取其 id，从 slidenotes 中获取相应描述
        调用 _generate_img 函数，生成对应图片，并进行图片替换
        最后将在 text 的最后追加 -g，表示图片已生成
        args:
            slide: python-pptx 中的 slide
        '''
        slide_notes_str = slide.notes_slide.notes_text_frame.text
        await self.replace_shapes_img(slide.shapes, slide_notes_str)

    # 批量替换图片，yield slide 的 id
    async def batch_replace_multimedia(self):
        '''
        批量替换图片，yield slide 的 id
        '''
        for slide in self.ppt.slides:
            await self.replace_slide_img(slide)
            self.replace_icons(slide.shapes)
            yield slide.slide_id

# =============
# 生成最终清洁版 ppt
# =============

    def clear_shapes(self, shapes):
        '''
        清除 shapes 中 @img-* 或 @svg-* 开头的 shape
        '''
        for shp in shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.clear_shapes(shp.shapes)
            elif hasattr(shp, 'has_text_frame') and shp.has_text_frame:
                text = shp.text_frame.text
                if text.startswith('@img') or text.startswith('@svg'):
                    shapes._spTree.remove(shp.element)


    def clear_slide(self, slide):
        '''
        清除 slide 中 @img-* 或 @svg-* 开头的 shape
        清除 slide_notes 中的 @img-*  至 @endimg 或 @svg-* 至 @endsvg 的内容
        '''
        slide_notes = slide.notes_slide.notes_text_frame.text
        slide_notes = re.sub(r'@img-[\s\S]*?@endimg', '', slide_notes)
        slide_notes = re.sub(r'@svg-[\s\S]*?@endsvg', '', slide_notes)
        slide.notes_slide.notes_text_frame.text = slide_notes

        self.clear_shapes(slide.shapes)

    def batch_clear_slide(self):
        '''
        批量清除 slide 中 @img-* 或 @svg-* 开头的 shape
        '''
        for slide in self.ppt.slides:
            self.clear_slide(slide)
